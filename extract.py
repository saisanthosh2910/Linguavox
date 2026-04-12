import requests
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

API_KEY = "K84999071788957"

# -------- PDF --------
def extract_text_from_pdf(file):
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

# -------- TXT --------
def extract_text_from_txt(file):
    return file.read().decode("utf-8")

# -------- DOCX --------
def extract_text_from_docx(file):
    doc = Document(file)
    return "\n".join([p.text for p in doc.paragraphs])

# -------- PPTX --------
def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# -------- IMAGE OCR --------
def extract_text_from_image(file):
    try:
        response = requests.post(
            "https://api.ocr.space/parse/image",
            files={"file": ("image.jpg", file.getvalue())},
            data={"apikey": API_KEY}
        )

        result = response.json()

        if result.get("IsErroredOnProcessing"):
            return "OCR Error"

        return result["ParsedResults"][0]["ParsedText"]

    except Exception as e:
        return str(e)